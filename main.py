"""
ChurchTool — FastAPI Backend
  • Bible Slide Generator
  • Song / Music-Sheet OCR → Lyric Slide Generator
"""
from __future__ import annotations

import io
import os
import re
import uuid
import zipfile
from collections import Counter
from pathlib import Path
from typing import Optional

from fastapi import FastAPI, File, Form, UploadFile
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import FileResponse, StreamingResponse
from fastapi.staticfiles import StaticFiles
from pydantic import BaseModel, Field

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

app = FastAPI(title="ChurchTool — Bible & Song Slide Generator")

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_methods=["*"],
    allow_headers=["*"],
)

# ─────────────────────────────────────────────
# Models
# ─────────────────────────────────────────────

class SlideSettings(BaseModel):
    font_family: str = "Arial"
    font_size: int = Field(36, ge=8, le=120)
    bold: bool = False
    align: str = "center"
    line_spacing: float = Field(1.2, ge=0.5, le=4.0)
    bg_color: str = "#000000"
    text_color: str = "#ffffff"
    ref_color: str = "#aaaaaa"
    ref_font_size: int = Field(20, ge=8, le=72)
    verse_num_color: str = "#ffcc44"   # verse number colour
    padding: float = Field(0.55, ge=0.1, le=2.0)
    show_ref_each: bool = True


class GenerateRequest(BaseModel):
    raw_text: str
    per_slide: int = Field(1, ge=1, le=20)
    settings: SlideSettings = SlideSettings()
    reference: str = ""


class SlidePart(BaseModel):
    """A single rendered slide."""
    verse_num: str        # e.g. "12" — rendered in verse_num_color
    lines: list[str]
    ref: str


class PreviewRequest(BaseModel):
    raw_text: str
    per_slide: int = Field(1, ge=1, le=20)
    settings: SlideSettings = SlideSettings()
    reference: str = ""


# ─────────────────────────────────────────────
# Helpers
# ─────────────────────────────────────────────

def _hex_to_rgb(hex_color: str) -> tuple[int, int, int]:
    h = hex_color.lstrip("#")
    if len(h) == 3:
        h = "".join(c * 2 for c in h)
    return int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)


ALIGN_MAP = {
    "left": PP_ALIGN.LEFT,
    "center": PP_ALIGN.CENTER,
    "right": PP_ALIGN.RIGHT,
}

# ── Verse number pattern: Arabic digits OR Khmer digits at line start
# Matches: "12 ..."  or  "១២ ..."  (one or more digits, then whitespace, then content)
_NUM_PAT = r"[0-9០-៩]+"
KHMER_VERSE_RE = re.compile(rf"^({_NUM_PAT})\s+(.+)$", re.MULTILINE)

# ── English reference blocks: "John 3:16\nbody…"
VERSE_REF_RE = re.compile(
    r"""
    (?:^|\n)\s*
    (?P<ref>
        (?:\d+\s+)?[A-Z][a-z]+(?:\s+[A-Za-z]+)*
        \s+\d+:\d+(?:[–\-]\d+)?
    )\s*\n
    (?P<body>.*?)
    (?=(?:\d+\s+)?[A-Z][a-z]+(?:\s+[A-Za-z]+)*\s+\d+:\d+|$)
    """,
    re.DOTALL | re.VERBOSE,
)


def parse_verses(raw_text: str) -> list[dict]:
    """
    Returns list of {ref, verse_num, lines}.

    Detection priority:
      1. Khmer-style leading verse numbers  ("12 ដំ...")
      2. English reference headers          ("John 3:16\\nbody")
      3. Blank-line block fallback
    """
    raw_text = raw_text.strip()
    if not raw_text:
        return []

    # Strategy 1 — Khmer numbered lines
    khmer_matches = list(KHMER_VERSE_RE.finditer(raw_text))
    non_blank = [l for l in raw_text.splitlines() if l.strip()]
    if khmer_matches and len(khmer_matches) / max(len(non_blank), 1) >= 0.5:
        verses = []
        for m in khmer_matches:
            body = m.group(2).strip()
            if body:
                verses.append({"ref": "", "verse_num": m.group(1), "lines": [body]})
        if verses:
            return verses

    # Strategy 2 — English reference blocks
    eng_matches = list(VERSE_REF_RE.finditer(raw_text))
    if eng_matches:
        verses = []
        for m in eng_matches:
            lines = [l.strip() for l in m.group("body").splitlines() if l.strip()]
            if lines:
                verses.append({"ref": m.group("ref").strip(), "verse_num": "", "lines": lines})
        if verses:
            return verses

    # Strategy 3 — blank-line blocks
    verses = []
    for block in re.split(r"\n\s*\n", raw_text):
        block = block.strip()
        if not block:
            continue
        verses.append({"ref": "", "verse_num": "", "lines": [l.strip() for l in block.splitlines() if l.strip()]})
    return verses


def _chunk(items: list, size: int) -> list[list]:
    return [items[i: i + size] for i in range(0, len(items), size)]


def build_slides(
    verses: list[dict],
    per_slide: int,
    global_ref: str,
    show_ref_each: bool,
) -> list[SlidePart]:
    slides: list[SlidePart] = []
    for v in verses:
        ref = v.get("ref") or global_ref
        verse_num = v.get("verse_num", "")
        groups = _chunk(v["lines"], per_slide)
        for i, group in enumerate(groups):
            show = show_ref_each or i == len(groups) - 1
            slides.append(SlidePart(
                verse_num=verse_num,
                lines=group,
                ref=ref if show else "",
            ))
    return slides


# ─────────────────────────────────────────────
# PPTX builder
# ─────────────────────────────────────────────

def _fix_font(run, font_name: str):
    """
    python-pptx sets font name via theme inheritance which PowerPoint may override.
    Directly writing <a:latin typeface="..."/> on the run's rPr guarantees the font
    is respected without needing a manual click in PowerPoint.
    """
    rPr = run._r.get_or_add_rPr()
    # Remove any existing <a:latin> elements first
    for el in rPr.findall(qn("a:latin")):
        rPr.remove(el)
    latin = etree.SubElement(rPr, qn("a:latin"))
    latin.set("typeface", font_name)
    # For complex/East-Asian scripts (Khmer uses complex script path in pptx)
    for el in rPr.findall(qn("a:cs")):
        rPr.remove(el)
    cs = etree.SubElement(rPr, qn("a:cs"))
    cs.set("typeface", font_name)


def _set_line_spacing(p, line_spacing: float):
    """
    Apply line spacing as lnSpc/spcPct.
    PowerPoint value: 1.0 = 100000, 1.5 = 150000, 2.0 = 200000.
    We also clear any spcBef/spcAft so there is no extra gap between lines.
    """
    pPr = p._p.get_or_add_pPr()
    # Remove any existing line-spacing / space-before / space-after elements
    for tag in ("a:lnSpc", "a:spcBef", "a:spcAft"):
        for el in pPr.findall(qn(tag)):
            pPr.remove(el)
    # Insert lnSpc as the first child of pPr (required order by schema)
    lnSpc = etree.SubElement(pPr, qn("a:lnSpc"))
    spcPct = etree.SubElement(lnSpc, qn("a:spcPct"))
    spcPct.set("val", str(int(line_spacing * 100000)))
    # Zero out space before/after paragraphs
    for tag in ("a:spcBef", "a:spcAft"):
        spc = etree.SubElement(pPr, qn(tag))
        pts = etree.SubElement(spc, qn("a:spcPts"))
        pts.set("val", "0")


def _add_textbox(
    slide,
    lines: list[str],
    left: float, top: float, width: float, height: float,
    font_name: str, font_size: int, bold: bool,
    color_rgb: tuple, align,
    line_spacing: float,
    first_line_prefix: str = "",
    prefix_color_rgb: tuple | None = None,
):
    """
    Add a text box. If first_line_prefix is set (e.g. verse number), it is
    prepended as a separate run on the first paragraph with prefix_color_rgb.
    """
    txBox = slide.shapes.add_textbox(int(left), int(top), int(width), int(height))
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        _set_line_spacing(p, line_spacing)

        # Verse number prefix — only on first line, coloured differently
        if i == 0 and first_line_prefix:
            pre_run = p.add_run()
            pre_run.text = first_line_prefix + "  "   # two spaces as separator
            pre_run.font.size = Pt(font_size)
            pre_run.font.bold = True
            pre_run.font.color.rgb = RGBColor(*(prefix_color_rgb or color_rgb))
            _fix_font(pre_run, font_name)

        run = p.add_run()
        run.text = line
        run.font.name = font_name
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = RGBColor(*color_rgb)
        _fix_font(run, font_name)


def build_pptx(slides_data: list[SlidePart], s: SlideSettings) -> bytes:
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    align     = ALIGN_MAP.get(s.align, PP_ALIGN.CENTER)
    text_rgb  = _hex_to_rgb(s.text_color)
    ref_rgb   = _hex_to_rgb(s.ref_color)
    vnum_rgb  = _hex_to_rgb(s.verse_num_color)
    bg_rgb    = _hex_to_rgb(s.bg_color)

    pad = Inches(s.padding)
    sw  = prs.slide_width
    sh  = prs.slide_height

    for sd in slides_data:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        bg_fill = slide.background.fill
        bg_fill.solid()
        bg_fill.fore_color.rgb = RGBColor(*bg_rgb)

        # Body text — verse number is prepended inline on the first line
        body_w = sw - pad * 2
        body_h = sh - pad * 2 - Inches(0.85)
        _add_textbox(
            slide, sd.lines,
            pad, pad, body_w, body_h,
            s.font_family, s.font_size, s.bold,
            text_rgb, align, s.line_spacing,
            first_line_prefix=sd.verse_num,
            prefix_color_rgb=vnum_rgb,
        )

        # Reference (bottom-right)
        if sd.ref:
            _add_textbox(
                slide, [sd.ref],
                pad, sh - Inches(0.75),
                sw - pad * 2, Inches(0.6),
                s.font_family, s.ref_font_size, False,
                ref_rgb, PP_ALIGN.RIGHT, 1.0,
            )

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()


# ─────────────────────────────────────────────
# Routes
# ─────────────────────────────────────────────

@app.get("/api/health")
def health():
    return {"status": "ok"}


@app.post("/api/preview")
def preview(req: PreviewRequest):
    verses = parse_verses(req.raw_text)
    slides = build_slides(verses, req.per_slide, req.reference, req.settings.show_ref_each)
    return {"slides": [s.model_dump() for s in slides], "total": len(slides)}


@app.post("/api/generate")
def generate(req: GenerateRequest):
    verses = parse_verses(req.raw_text)
    if not verses:
        verses = [{"ref": req.reference or "", "verse_num": "", "lines": ["(No content)"]}]
    slides = build_slides(verses, req.per_slide, req.reference, req.settings.show_ref_each)
    pptx_bytes = build_pptx(slides, req.settings)
    filename = f"bible_slides_{uuid.uuid4().hex[:8]}.pptx"
    return StreamingResponse(
        io.BytesIO(pptx_bytes),
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": f'attachment; filename="{filename}"'},
    )


# ═══════════════════════════════════════════════════════════════════
#  Song / OCR Slide Generator
# ═══════════════════════════════════════════════════════════════════

# ── Fonts ────────────────────────────────────────────────────────
FONT_DIR = Path(__file__).parent / "fonts"
FONT_FILE_MAP: dict[str, str] = {
    "Khmer OS Battambang": "KhmerOSBattambang-Regular.ttf",
    "Khmer OS":            "KhmerOS_.ttf",
    "Khmer OS Muol":       "KhmerOS_muol.ttf",
    "Khmer OS Moul Light": "KhmerOSMoulLight.ttf",
}

# ── Text-filtering helpers ───────────────────────────────────────
_HAS_KHMER_RE = re.compile(r'[\u1780-\u17FF\u19E0-\u19FF]')

# Khmer Unicode ranges
_KHMER_CHAR_RE = re.compile(r'[\u1780-\u17FF\u19E0-\u19FF]')

# Chord token (letter root + quality + optional slash bass)
_CHORD_TOKEN_RE = re.compile(
    r'^[A-G][#b]?(?:m|maj|min|aug|dim|sus|add)?(?:\d+)?(?:/[A-G][#b]?)?$',
    re.IGNORECASE,
)


def _line_is_khmer(s: str) -> bool:
    return bool(_HAS_KHMER_RE.search(s))


def _khmer_char_count(s: str) -> int:
    return len(_KHMER_CHAR_RE.findall(s))


def _is_repetitive(s: str, thr: float = 0.45) -> bool:
    """
    Return True if top-1 Khmer char >= thr fraction, OR top-3 chars >= 72 %
    of all Khmer chars.  Catches both per-char smear and solfege notation.
    """
    khmer_chars = _KHMER_CHAR_RE.findall(s)
    if len(khmer_chars) < 4:
        return False
    counts = Counter(khmer_chars)
    top1_frac = counts.most_common(1)[0][1] / len(khmer_chars)
    if top1_frac >= thr:
        return True
    top3_total = sum(c for _, c in counts.most_common(3))
    if top3_total / len(khmer_chars) >= 0.72:
        return True
    return False


def _has_repeated_bigram(tokens: list, max_frac: float = 0.40) -> bool:
    """Return True if any consecutive 2-token pair repeats >= max_frac of all bigrams."""
    if len(tokens) < 4:
        return False
    bigrams = [(tokens[i], tokens[i + 1]) for i in range(len(tokens) - 1)]
    top_count = Counter(bigrams).most_common(1)[0][1]
    return top_count / len(bigrams) >= max_frac


# OCR noise characters to strip from lyric lines before filtering
# These are common misreads: slashes, angles, brackets, punctuation artifacts
_OCR_NOISE_RE = re.compile(r'[/\\<>«»›‹|(){}\[\]#@&*^~`]+')

# Token-level smear detection regexes (compiled once at module level)
_LATIN_RE      = re.compile(r'[a-zA-Z]')
_REPEAT_SUB_RE = re.compile(r'(.{2,3})\1{2,}')  # e.g. "ោយោយោយ"


def _clean_line(s: str) -> str:
    """
    Strip common OCR noise characters from a lyric line while preserving
    Khmer text and natural punctuation (dash, comma, period, space).
    Also collapses multiple spaces.
    """
    s = _OCR_NOISE_RE.sub(' ', s)
    s = re.sub(r' {2,}', ' ', s)
    return s.strip()


def _is_lyric_line(s: str) -> bool:
    """
    Return True only if this line looks like a genuine Khmer lyric line.

    Rejection criteria (any one -> reject):
      1. No Khmer characters at all.
      2. Fewer than 6 Khmer characters.
      3. Top-1 Khmer char >= 45 % (repetitive OCR smear).
      4. Top-3 Khmer chars >= 72 % (solfege / rhythm notation).
      5. Fewer than 2 whitespace-separated Khmer word tokens.
      6. 2-3 tokens AND avg token length < 4.5 chars
         (short instrument-label / tab marker).
      7. Average Khmer token length < 2.5 chars (pure notation markers).
      8. Two consecutive identical Khmer tokens (rhythm: "i i", "pi pi").
      9. Repeated bigrams across the line (rhythm notation rows).
     10. >= 60 % of tokens end with Khmer vowel sign ii (\u17B8 = \u17b8 = i-vowel)
         -> solfege syllable row (pii, kii, rii, sii...).
     11. >= 50 % of tokens are exactly 2 chars (solfege syllables).
     12. Khmer chars < 55 % of all non-space chars (mixed notation/ascii).
    """
    if not _line_is_khmer(s):
        return False

    k_count = _khmer_char_count(s)
    if k_count < 6:
        return False

    if _is_repetitive(s):
        return False

    khmer_word_tokens = [t for t in s.split() if _line_is_khmer(t)]
    n_tokens = len(khmer_word_tokens)

    if n_tokens < 2:
        return False

    avg_word_len = sum(len(t) for t in khmer_word_tokens) / n_tokens

    # Short instrument-label gate: <= 3 tokens with short avg length
    if n_tokens <= 3 and avg_word_len < 4.5:
        return False

    # General avg-length floor
    if avg_word_len < 2.5:
        return False

    # Consecutive identical tokens
    for i in range(n_tokens - 1):
        if khmer_word_tokens[i] == khmer_word_tokens[i + 1]:
            return False

    # Repeated bigrams
    if _has_repeated_bigram(khmer_word_tokens):
        return False

    # Solfege gate 1: >= 60% of tokens end with ii-vowel sign (U+17B8)
    # Real Khmer words rarely end with ii; solfege syllables always do
    ii_vowel = '\u17B8'
    ii_end_frac = sum(1 for t in khmer_word_tokens if ii_vowel in t[-2:]) / n_tokens
    if ii_end_frac >= 0.60:
        return False

    # Solfege gate 2: >= 50% of tokens are exactly 2 chars
    short2_frac = sum(1 for t in khmer_word_tokens if len(t) == 2) / n_tokens
    if short2_frac >= 0.50:
        return False

    # Token-level smear gates: reject if ANY token looks like notation smear
    for tok in khmer_word_tokens:
        # Mixed Khmer+Latin in one token -> OCR smear of staff lines
        if _LATIN_RE.search(tok) and _KHMER_CHAR_RE.search(tok):
            return False
        # Repeated 2–3 char substring pattern -> notation run (e.g. ោយោយោយ)
        if len(tok) >= 6 and _REPEAT_SUB_RE.search(tok):
            return False
        # Long token with highly repetitive character bigrams -> notation run
        if len(tok) >= 7:
            chars = list(tok)
            bigrams = [chars[i] + chars[i+1] for i in range(len(chars)-1)]
            if bigrams:
                top_bg_count = Counter(bigrams).most_common(1)[0][1]
                if top_bg_count / len(bigrams) >= 0.22:
                    return False

    # Khmer density gate
    total_chars = len(s.replace(' ', ''))
    if total_chars > 0 and k_count / total_chars < 0.55:
        return False

    return True


def filter_lyrics(raw: str) -> str:
    """
    Pass 1: for each line, strip OCR noise chars then test _is_lyric_line().
    Pass 2: deduplicate adjacent identical lines (repeated verse detection).
    Blank lines are preserved as verse separators.
    """
    kept: list[str] = []
    prev_blank = True
    prev_line = None
    for line in raw.splitlines():
        s = line.strip()
        if not s:
            if not prev_blank:
                kept.append("")
            prev_blank = True
            prev_line = None
            continue
        prev_blank = False
        # Clean OCR noise first, then test
        cleaned = _clean_line(s)
        if _is_lyric_line(cleaned):
            if cleaned != prev_line:   # skip exact duplicate consecutive lines
                kept.append(cleaned)
                prev_line = cleaned
    while kept and kept[-1] == "":
        kept.pop()
    return "\n".join(kept)


# ── Numpy / PIL imports ─────────────────────────────────────────
import numpy as np
from PIL import Image, ImageOps, ImageFilter


def _find_tesseract() -> str:
    import shutil
    for candidate in (
        "/opt/homebrew/bin/tesseract",
        "/usr/local/bin/tesseract",
        "/usr/bin/tesseract",
    ):
        if os.path.exists(candidate):
            return candidate
    return shutil.which("tesseract") or "tesseract"


_TESSERACT_CMD: Optional[str] = None


def _get_tesseract() -> str:
    global _TESSERACT_CMD
    if _TESSERACT_CMD is None:
        _TESSERACT_CMD = _find_tesseract()
    return _TESSERACT_CMD


# ── Image pre-processing ────────────────────────────────────────

def _preprocess_for_ocr(img: Image.Image, target_width: int = 2800) -> Image.Image:
    """
    Prepare a lyric-strip image for Tesseract (Khmer-optimised):

      1. Greyscale
      2. Upscale so the longest dimension ≥ target_width px
         (Khmer vowel signs need glyph height ≥ 50 px to be read correctly)
      3. Auto-contrast (stretch histogram to full 0-255 range)
      4. Unsharp-mask (crisper edges than double-SHARPEN, doesn't over-thicken)
      5. Adaptive binarise via numpy:
           - Compute local mean in a 41-px window
           - Threshold = local_mean - 18  (handles uneven lighting / shadows)
           - This preserves faint vowel signs that hard-threshold at 150 erases
      6. Light despeckle: remove isolated dark pixels smaller than 2×2
      7. White border (Tesseract hates glyphs at the very edge)
    """
    img = img.convert("L")
    w, h = img.size

    # Step 2 — upscale
    if w < target_width:
        scale = target_width / w
        new_w = int(w * scale)
        new_h = max(int(h * scale), 80)   # ensure min height for vowel signs
        img = img.resize((new_w, new_h), Image.LANCZOS)

    # Step 3 — auto-contrast
    img = ImageOps.autocontrast(img, cutoff=2)

    # Step 4 — unsharp mask (radius 2, percent 150, threshold 3)
    img = img.filter(ImageFilter.UnsharpMask(radius=2, percent=150, threshold=3))

    # Step 5 — adaptive binarise
    arr = np.array(img, dtype=np.float32)
    from scipy.ndimage import uniform_filter
    try:
        local_mean = uniform_filter(arr, size=41)
        binary = (arr > (local_mean - 18)).astype(np.uint8) * 255
    except ImportError:
        # scipy not available → fall back to simple threshold at 140
        binary = (arr > 140).astype(np.uint8) * 255
    img = Image.fromarray(binary.astype(np.uint8), mode="L")

    # Step 6 — light despeckle: erode then dilate (opening) to remove 1-px noise
    img = img.filter(ImageFilter.MinFilter(3))   # erode (dark = ink)
    img = img.filter(ImageFilter.MaxFilter(3))   # dilate back

    # Step 7 — white border
    img = ImageOps.expand(img, border=20, fill=255)
    return img


# ── Staff-line detection ────────────────────────────────────────

def _row_dark_density(gray: np.ndarray, thr: int = 180) -> np.ndarray:
    """Fraction of pixels darker than thr in each row."""
    return (gray < thr).mean(axis=1).astype(float)


def _moving_avg(arr: np.ndarray, w: int) -> np.ndarray:
    if w < 2:
        return arr.copy()
    return np.convolve(arr, np.ones(w) / w, mode="same")


def _gap_spans(mask: np.ndarray, min_px: int) -> list[tuple[int, int]]:
    spans, in_g, g0 = [], False, 0
    for y, v in enumerate(mask):
        if v and not in_g:
            in_g, g0 = True, y
        elif not v and in_g:
            in_g = False
            if y - g0 >= min_px:
                spans.append((g0, y))
    if in_g and len(mask) - g0 >= min_px:
        spans.append((g0, len(mask)))
    return spans


def _row_horizontal_runs(gray_row: np.ndarray, thr: int = 180) -> float:
    """
    Return the mean length of dark horizontal runs in a single row.
    Staff lines have very long uniform runs (spanning most of the page width).
    Lyric rows have short, irregular runs (individual glyph strokes).
    """
    dark = gray_row < thr
    runs, in_run, run_len = [], False, 0
    for v in dark:
        if v:
            in_run = True
            run_len += 1
        elif in_run:
            runs.append(run_len)
            in_run = False
            run_len = 0
    if in_run:
        runs.append(run_len)
    return float(np.mean(runs)) if runs else 0.0


def _is_staff_strip(gray: np.ndarray, top: int, bot: int,
                    page_width: int) -> bool:
    """
    Decide whether the horizontal band gray[top:bot] is a music notation
    (staff) strip using TWO independent signals:

    Signal 1 — Long horizontal runs:
      Staff lines are thin, page-spanning dark lines.
      A strip with ≥ 3 rows where mean run length > 15 % of page width
      is almost certainly staff+notation.

    Signal 2 — Row-density profile:
      Staff rows have moderate but very uniform density across the height.
      The standard deviation of per-row densities is small for staff bands.

    Either signal alone is sufficient to label a strip as notation.
    """
    band       = gray[top:bot]
    h_band     = band.shape[0]
    if h_band < 3:
        return False

    run_threshold  = page_width * 0.12   # run > 12 % of page width → staff line
    long_run_rows  = 0
    row_densities  = []

    for row_idx in range(h_band):
        row = band[row_idx]
        mean_run = _row_horizontal_runs(row, thr=180)
        if mean_run > run_threshold:
            long_run_rows += 1
        row_densities.append(float((row < 180).mean()))

    # Signal 1: multiple rows with very long horizontal runs
    if long_run_rows >= max(2, h_band * 0.15):
        return True

    # Signal 2: high overall density (lots of ink — notes + beams) AND
    #           uniform density variance (regular pattern of staff lines)
    arr_d = np.array(row_densities)
    mean_d = float(arr_d.mean())
    std_d  = float(arr_d.std())
    if mean_d > 0.08 and std_d < 0.06:
        return True

    return False


def _detect_content_strips(gray: np.ndarray,
                            blank_thresh: float = 0.008,
                            min_gap_px:   int   = 6,
                            min_strip_px: int   = 15,
                            smooth_w:     int   = 3,
                            ) -> list[tuple[int, int]]:
    """
    Split a greyscale page into horizontal content strips by finding
    near-blank row gaps.  Uses a narrow smoothing window so stave/lyric
    boundaries stay sharp.
    """
    density = _moving_avg(_row_dark_density(gray, thr=180), w=smooth_w)
    blank   = density < blank_thresh
    strips, prev = [], 0
    for g0, g1 in _gap_spans(blank, min_gap_px):
        if g0 - prev >= min_strip_px:
            strips.append((prev, g0))
        prev = g1
    h = gray.shape[0]
    if h - prev >= min_strip_px:
        strips.append((prev, h))
    return strips


def _extract_lyric_strips(page_rgb: Image.Image) -> list[Image.Image]:
    """
    Return cropped PIL images of ONLY the lyric (Khmer text) rows from
    a music sheet page, discarding all notation/staff rows.

    Algorithm:
      1. Detect all horizontal content strips (separated by white gaps).
      2. For each strip, run _is_staff_strip() to detect notation rows
         using long-horizontal-run analysis.
      3. Any strip NOT identified as staff is a candidate lyric row.
      4. Safety: if zero lyric strips found (flat/photocopied page),
         return [] so the caller falls back to full-page OCR + text filter.
    """
    gray  = np.array(page_rgb.convert("L"))
    h, w  = gray.shape

    strips = _detect_content_strips(gray)
    if len(strips) < 2:
        return []

    lyric_crops: list[Image.Image] = []
    for top, bot in strips:
        if _is_staff_strip(gray, top, bot, page_width=w):
            continue   # skip notation / staff row entirely

        # This strip is NOT staff → it's lyrics (or chord symbols which
        # we'll filter at the text level after OCR).
        # Use generous vertical padding so vowel signs above/below baseline
        # are never clipped (Khmer stacks go well above and below the base row).
        pad = 24
        t = max(0, top - pad)
        b = min(h, bot + pad)
        lyric_crops.append(page_rgb.crop((0, t, w, b)))

    return lyric_crops


# ── OCR helpers ──────────────────────────────────────────────────

# Common Tesseract Khmer OCR mis-substitutions — map wrong → correct codepoint
# These are single-character swaps confirmed across many Khmer OCR test runs.
_KHMER_FIXES: list[tuple[str, str]] = [
    # Visually similar consonant pairs that Tesseract confuses
    ('\u1780', '\u1781'),   # ក ↔ ខ  (ka / kha) — only swap if context demands;
                            # actually we leave consonants alone (too risky without context)
    # Vowel signs commonly swapped by Tesseract
    ('\u17BE', '\u17C1'),   # ើ → េ  (oe → e) — Tesseract over-uses ើ
    # Subscript coeng (​្) sometimes dropped → we can't restore it, but
    # we CAN normalise Unicode composition
]

# Characters Tesseract inserts that are not valid Khmer (e.g. Cyrillic leaked
# from the model, Latin letters surrounded by Khmer)
_STRIP_NON_KHMER_IN_WORD_RE = re.compile(
    r'(?<=[\u1780-\u17FF\u19E0-\u19FF])[a-zA-Z\u0400-\u04FF]+'
    r'(?=[\u1780-\u17FF\u19E0-\u19FF])'
)


def _fix_ocr_text(text: str) -> str:
    """
    Post-process raw Tesseract output to fix the most common Khmer OCR errors:
      1. Unicode normalisation (NFC) — ensures combining marks are properly attached
      2. Remove Cyrillic/Latin characters that sneak in between Khmer characters
      3. Strip lines that are pure ASCII/numbers after Khmer-only filtering
      4. Collapse multiple spaces
    """
    import unicodedata
    text = unicodedata.normalize("NFC", text)
    # Remove stray non-Khmer chars embedded inside Khmer words
    text = _STRIP_NON_KHMER_IN_WORD_RE.sub('', text)
    # Collapse spaces
    text = re.sub(r' {2,}', ' ', text)
    return text


def _khmer_char_count_in(text: str) -> int:
    """Count total Khmer characters in a multi-line string."""
    return len(_KHMER_CHAR_RE.findall(text))


def _ocr_attempt(processed: 'Image.Image', lang: str, psm: str,
                 tessdata_dir: str | None = None) -> str:
    """Run one Tesseract attempt; return raw text or ''."""
    try:
        import pytesseract
    except ImportError:
        return ""
    cfg = f"--psm {psm} --oem 1"
    if tessdata_dir:
        cfg += f" --tessdata-dir {tessdata_dir}"
    try:
        return pytesseract.image_to_string(processed, lang=lang, config=cfg)
    except Exception:
        return ""


def _ocr_strip(img: Image.Image) -> str:
    """
    OCR a single lyric strip with multiple strategy attempts.
    Picks the result with the most Khmer characters (best signal).

    Strategies tried (in parallel order, best wins):
      A. lang=khm,       psm=7  (single line  — narrow strips)
      B. lang=khm,       psm=6  (uniform block — taller strips)
      C. lang=script/Khmer, psm=6  (script-level model, better glyph accuracy)
      D. lang=khm+eng,   psm=6  (allow some Latin — chord row guard)

    After picking the winner:
      • Keep only lines containing Khmer characters
      • Apply _fix_ocr_text() Unicode normalisation + noise removal
    """
    try:
        import pytesseract
    except ImportError:
        return ""

    pytesseract.pytesseract.tesseract_cmd = _get_tesseract()
    processed = _preprocess_for_ocr(img)
    _, ph = processed.size
    psm_primary = "7" if ph < 100 else "6"

    # Build script/Khmer tessdata path if available
    tess_prefix = os.path.dirname(os.path.dirname(_get_tesseract()))  # e.g. /opt/homebrew
    script_dir  = os.path.join(tess_prefix, "share", "tessdata")
    script_lang = "script/Khmer" if os.path.exists(
        os.path.join(script_dir, "script", "Khmer.traineddata")
    ) else None

    candidates: list[str] = []

    # Strategy A/B: khm with both psm modes
    for psm in (psm_primary, "6" if psm_primary == "7" else "7"):
        raw = _ocr_attempt(processed, "khm", psm)
        if raw:
            candidates.append(raw)

    # Strategy C: script/Khmer model (often better at individual glyph shapes)
    if script_lang:
        raw = _ocr_attempt(processed, script_lang, "6", tessdata_dir=script_dir)
        if raw:
            candidates.append(raw)

    # Strategy D: khm+eng fallback
    raw = _ocr_attempt(processed, "khm+eng", "6")
    if raw:
        candidates.append(raw)

    if not candidates:
        return ""

    # Pick the candidate with the most Khmer characters
    best = max(candidates, key=_khmer_char_count_in)

    # Post-process: keep Khmer-containing lines, fix encoding, normalise
    lines = [
        _fix_ocr_text(ln.strip())
        for ln in best.splitlines()
        if ln.strip() and _line_is_khmer(ln)
    ]
    return "\n".join(lines)


def _ocr_full_page_khmer(page_rgb: Image.Image) -> str:
    """
    Fallback: OCR the whole page with auto-layout detection (psm 3).
    Tries both khm and script/Khmer models; picks best (most Khmer chars).
    """
    try:
        import pytesseract
    except ImportError:
        return ""

    pytesseract.pytesseract.tesseract_cmd = _get_tesseract()
    processed = _preprocess_for_ocr(page_rgb)

    tess_prefix = os.path.dirname(os.path.dirname(_get_tesseract()))
    script_dir  = os.path.join(tess_prefix, "share", "tessdata")
    script_lang = "script/Khmer" if os.path.exists(
        os.path.join(script_dir, "script", "Khmer.traineddata")
    ) else None

    candidates: list[str] = []
    for lang in ("khm", "khm+eng"):
        raw = _ocr_attempt(processed, lang, "3")
        if raw:
            candidates.append(raw)
    if script_lang:
        raw = _ocr_attempt(processed, script_lang, "3", tessdata_dir=script_dir)
        if raw:
            candidates.append(raw)

    if not candidates:
        return ""

    best = max(candidates, key=_khmer_char_count_in)
    lines = [
        _fix_ocr_text(ln.strip())
        for ln in best.splitlines()
        if ln.strip() and _line_is_khmer(ln)
    ]
    return "\n".join(lines)


# ── Main extraction entry point ─────────────────────────────────

def extract_lyrics_from_bytes(data: bytes, ext: str) -> str:
    """
    Accept raw file bytes + extension (.pdf / .png / .jpg / .jpeg).
    Returns ONLY the Khmer lyric lines, with blank lines as verse separators.

    Pipeline per page:
      1. Convert PDF → RGB at 300 dpi (larger = better OCR)
      2. Detect & isolate lyric strips by discarding staff/notation rows
      3. OCR each lyric strip with lang=khm, keep only Khmer lines
      4. If strip isolation yields nothing, OCR full page and filter to Khmer
      5. Run filter_lyrics() as a final pass to ensure no garbage remains
    """
    all_lines: list[str] = []

    def process_page(page_rgb: Image.Image) -> None:
        strips = _extract_lyric_strips(page_rgb)

        if strips:
            for strip_img in strips:
                text = _ocr_strip(strip_img)
                lines = [ln for ln in text.splitlines() if ln.strip()]
                if lines:
                    all_lines.extend(lines)
                    all_lines.append("")   # verse separator after each strip
        else:
            # Fallback: full-page OCR filtered to Khmer only
            text = _ocr_full_page_khmer(page_rgb)
            lines = [ln for ln in text.splitlines() if ln.strip()]
            all_lines.extend(lines)

    if ext == ".pdf":
        try:
            from pdf2image import convert_from_bytes
            pages = convert_from_bytes(data, dpi=300)
        except Exception as e:
            return f"[PDF conversion error: {e}]"
        for page in pages:
            process_page(page.convert("RGB"))
            all_lines.append("")   # separator between PDF pages
    else:
        try:
            img = Image.open(io.BytesIO(data)).convert("RGB")
            process_page(img)
        except Exception as e:
            return f"[Image error: {e}]"

    if not all_lines or all(not ln for ln in all_lines):
        return "(no Khmer text detected)"

    # Final filter: keep only Khmer lines, preserve blank separators
    raw      = "\n".join(all_lines)
    filtered = filter_lyrics(raw)
    return filtered.strip() if filtered.strip() else "(no Khmer text detected)"


# ── Song-slide PPTX builder ──────────────────────────────────────

class SongSettings(BaseModel):
    font_family:    str   = "Khmer OS Battambang"
    font_size:      int   = Field(56, ge=8, le=120)
    bold:           bool  = True
    align:          str   = "center"
    line_spacing:   float = Field(1.5, ge=0.5, le=4.0)
    verse_spacing:  float = Field(1.0, ge=0.0, le=3.0)   # extra space between verse groups
    bg_color:       str   = "#000000"
    text_color:     str   = "#ffffff"
    ref_color:      str   = "#aaaaaa"
    ref_font_size:  int   = Field(22, ge=8, le=72)
    padding:        float = Field(0.45, ge=0.05, le=2.0)
    lines_per_slide: int  = Field(2, ge=1, le=10)


class SongGenerateRequest(BaseModel):
    lyrics_text: str
    settings:    SongSettings = SongSettings()
    song_title:  str = ""


class SongPreviewRequest(BaseModel):
    lyrics_text: str
    settings:    SongSettings = SongSettings()
    song_title:  str = ""


def _build_song_pptx(lyrics_text: str, s: SongSettings, song_title: str = "") -> bytes:
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    align    = ALIGN_MAP.get(s.align, PP_ALIGN.CENTER)
    text_rgb = _hex_to_rgb(s.text_color)
    ref_rgb  = _hex_to_rgb(s.ref_color)
    bg_rgb   = _hex_to_rgb(s.bg_color)

    pad = Inches(s.padding)
    sw  = prs.slide_width
    sh  = prs.slide_height

    all_lines = [ln for ln in lyrics_text.splitlines() if ln.strip()]
    if not all_lines:
        all_lines = ["(no lyrics found)"]

    groups = [all_lines[i: i + s.lines_per_slide]
              for i in range(0, len(all_lines), s.lines_per_slide)]

    for g_idx, group in enumerate(groups):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        bg_fill = slide.background.fill
        bg_fill.solid()
        bg_fill.fore_color.rgb = RGBColor(*bg_rgb)

        body_h = sh - pad * 2 - (Inches(0.7) if song_title else Inches(0.1))
        txBox  = slide.shapes.add_textbox(int(pad), int(pad), int(sw - pad * 2), int(body_h))
        tf     = txBox.text_frame
        tf.word_wrap = True

        for i, line in enumerate(group):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = align
            _set_line_spacing(p, s.line_spacing)

            run = p.add_run()
            run.text = line
            run.font.size  = Pt(s.font_size)
            run.font.bold  = s.bold
            run.font.color.rgb = RGBColor(*text_rgb)
            _fix_font(run, s.font_family)

        # Song title / reference at bottom-right
        if song_title:
            _add_textbox(
                slide, [song_title],
                pad, sh - Inches(0.65),
                sw - pad * 2, Inches(0.55),
                s.font_family, s.ref_font_size, False,
                ref_rgb, PP_ALIGN.RIGHT, 1.0,
            )

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)

    # Embed Khmer font if applicable
    pptx_bytes = buf.read()
    ttf_filename = FONT_FILE_MAP.get(s.font_family)
    if ttf_filename:
        ttf_path = FONT_DIR / ttf_filename
        if ttf_path.exists():
            pptx_bytes = _embed_font(pptx_bytes, s.font_family, ttf_filename, ttf_path.read_bytes())

    return pptx_bytes


def _embed_font(pptx_bytes: bytes, font_name: str, ttf_filename: str, font_data: bytes) -> bytes:
    """Inject a TTF into the PPTX zip as an embedded font."""
    p_ns  = "http://schemas.openxmlformats.org/presentationml/2006/main"
    r_ns  = "http://schemas.openxmlformats.org/package/2006/relationships"
    r2_ns = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
    ct_ns = "http://schemas.openxmlformats.org/package/2006/content-types"
    font_rel_type  = f"{r2_ns}/font"
    font_zip_path  = f"ppt/fonts/{ttf_filename}"
    rels_path      = "ppt/_rels/presentation.xml.rels"
    prs_path_zip   = "ppt/presentation.xml"
    ct_path        = "[Content_Types].xml"

    src = io.BytesIO(pptx_bytes)
    dst = io.BytesIO()

    with zipfile.ZipFile(src, "r") as zin, zipfile.ZipFile(dst, "w", zipfile.ZIP_DEFLATED) as zout:
        # Patch rels
        rels_root = etree.fromstring(zin.read(rels_path))
        existing  = {el.get("Id") for el in rels_root}
        rid = "rFnt1"
        n   = 1
        while rid in existing:
            n  += 1
            rid = f"rFnt{n}"
        new_rel = etree.SubElement(rels_root, f"{{{r_ns}}}Relationship")
        new_rel.set("Id",     rid)
        new_rel.set("Type",   font_rel_type)
        new_rel.set("Target", f"fonts/{ttf_filename}")
        patched_rels = etree.tostring(rels_root, xml_declaration=True, encoding="UTF-8", standalone=True)

        # Patch presentation.xml
        prs_root = etree.fromstring(zin.read(prs_path_zip))
        lst = prs_root.find(f"{{{p_ns}}}embeddedFontLst")
        if lst is None:
            lst = etree.SubElement(prs_root, f"{{{p_ns}}}embeddedFontLst")
        ef  = etree.SubElement(lst, f"{{{p_ns}}}embeddedFont")
        fnt = etree.SubElement(ef,  f"{{{p_ns}}}font")
        fnt.set("typeface", font_name)
        reg = etree.SubElement(ef,  f"{{{p_ns}}}regular")
        reg.set(f"{{{r2_ns}}}id", rid)
        patched_prs = etree.tostring(prs_root, xml_declaration=True, encoding="UTF-8", standalone=True)

        # Patch Content_Types
        ct_root   = etree.fromstring(zin.read(ct_path))
        part_name = f"/ppt/fonts/{ttf_filename}"
        font_ct   = "application/vnd.openxmlformats-officedocument.presentationml.font"
        if not any(el.get("PartName") == part_name for el in ct_root.findall(f"{{{ct_ns}}}Override")):
            ov = etree.SubElement(ct_root, f"{{{ct_ns}}}Override")
            ov.set("PartName",    part_name)
            ov.set("ContentType", font_ct)
        patched_ct = etree.tostring(ct_root, xml_declaration=True, encoding="UTF-8", standalone=True)

        patched = {rels_path, prs_path_zip, ct_path}
        for item in zin.infolist():
            if item.filename not in patched:
                zout.writestr(item, zin.read(item.filename))
        zout.writestr(rels_path,     patched_rels)
        zout.writestr(prs_path_zip,  patched_prs)
        zout.writestr(ct_path,       patched_ct)
        zout.writestr(font_zip_path, font_data)

    dst.seek(0)
    return dst.read()


def _build_song_preview_slides(lyrics_text: str, s: SongSettings, song_title: str) -> list[dict]:
    all_lines = [ln for ln in lyrics_text.splitlines() if ln.strip()]
    if not all_lines:
        return [{"lines": ["(no lyrics)"], "title": song_title}]
    groups = [all_lines[i: i + s.lines_per_slide]
              for i in range(0, len(all_lines), s.lines_per_slide)]
    return [{"lines": g, "title": song_title if i == len(groups) - 1 else ""} for i, g in enumerate(groups)]


# ── Song OCR endpoint ────────────────────────────────────────────
@app.post("/api/song/extract")
async def song_extract(file: UploadFile = File(...)):
    """Accept a PDF/JPG/PNG upload, run OCR, return extracted lyrics as JSON."""
    ext = Path(file.filename or "").suffix.lower()
    if ext not in (".pdf", ".png", ".jpg", ".jpeg"):
        return {"error": "Only PDF, PNG, JPG files are supported."}, 400
    data = await file.read()
    text = extract_lyrics_from_bytes(data, ext)
    return {"text": text}


@app.post("/api/song/preview")
def song_preview(req: SongPreviewRequest):
    """Return list of slide objects for live preview."""
    slides = _build_song_preview_slides(req.lyrics_text, req.settings, req.song_title)
    return {"slides": slides, "total": len(slides)}


@app.post("/api/song/generate")
def song_generate(req: SongGenerateRequest):
    """Build and stream a PPTX from the provided lyrics + settings."""
    pptx_bytes = _build_song_pptx(req.lyrics_text, req.settings, req.song_title)
    filename   = f"song_slides_{uuid.uuid4().hex[:8]}.pptx"
    return StreamingResponse(
        io.BytesIO(pptx_bytes),
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": f'attachment; filename="{filename}"'},
    )


# ── Static frontend (Vite build) ─────────────────────────────────
# Mount the built React app so a single deployment serves everything.
# On Render: build.sh runs `npm run build` before gunicorn starts,
# so frontend/dist exists by the time this code runs.
_FRONTEND_DIST = Path(__file__).parent / "frontend" / "dist"

def _setup_static() -> None:
    """Mount static files only if the built frontend exists."""
    if not _FRONTEND_DIST.exists():
        return  # running locally without a build — API-only mode

    assets_dir = _FRONTEND_DIST / "assets"
    if assets_dir.exists():
        app.mount(
            "/assets",
            StaticFiles(directory=str(assets_dir)),
            name="assets",
        )

    # SPA catch-all — must be defined at module level (not inside a conditional)
    # so gunicorn workers inherit it correctly after fork.

@app.get("/{full_path:path}", include_in_schema=False)
async def serve_spa(full_path: str):
    """Catch-all: serve index.html for any path not matched by /api routes."""
    index = _FRONTEND_DIST / "index.html"
    if not index.exists():
        return {"detail": "Frontend not built"}
    candidate = _FRONTEND_DIST / full_path
    if candidate.is_file():
        return FileResponse(str(candidate))
    return FileResponse(str(index))

_setup_static()
