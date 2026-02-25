from flask import Flask, render_template, request, send_file, jsonify
from pdf2image import convert_from_path
from PIL import Image, ImageOps, ImageDraw, ImageFont
import numpy as np
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os, uuid, io, textwrap

app = Flask(__name__)
UPLOAD_FOLDER = "uploads"
OUTPUT_FOLDER = "output"
os.makedirs(UPLOAD_FOLDER, exist_ok=True)
os.makedirs(OUTPUT_FOLDER, exist_ok=True)


# ══════════════════════════════════════════════════════════════════
#  Shared PPTX helpers
# ══════════════════════════════════════════════════════════════════

def new_prs(bg_hex: str = "000000") -> Presentation:
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    return prs


def blank_slide(prs: Presentation, bg_hex: str = "000000"):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    r = int(bg_hex[0:2], 16)
    g = int(bg_hex[2:4], 16)
    b = int(bg_hex[4:6], 16)
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(r, g, b)
    return slide


def hex_to_rgb(h: str):
    h = h.lstrip("#")
    return int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)


def add_text_box(slide, text: str, left, top, width, height,
                 font_name: str, font_size: int,
                 bold: bool, color_hex: str, align: str = "center"):
    from pptx.util import Pt, Emu
    txBox = slide.shapes.add_textbox(int(left), int(top), int(width), int(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = {"left": PP_ALIGN.LEFT,
                   "center": PP_ALIGN.CENTER,
                   "right": PP_ALIGN.RIGHT}.get(align, PP_ALIGN.CENTER)
    run = p.runs[0]
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    r, g, b = hex_to_rgb(color_hex)
    run.font.color.rgb = RGBColor(r, g, b)


def save_prs(prs: Presentation, name: str) -> str:
    path = os.path.join(OUTPUT_FOLDER, name)
    prs.save(path)
    return path


# ══════════════════════════════════════════════════════════════════
#  Notation-cut helpers  (original logic preserved)
# ══════════════════════════════════════════════════════════════════

def row_dark_density(gray: np.ndarray, threshold: int = 240) -> np.ndarray:
    return (gray < threshold).mean(axis=1).astype(float)


def moving_avg(arr: np.ndarray, w: int = 9) -> np.ndarray:
    return np.convolve(arr, np.ones(w) / w, mode="same")


def gap_spans(is_blank: np.ndarray, min_px: int) -> list:
    spans, in_g, g0 = [], False, 0
    for y, b in enumerate(is_blank):
        if b and not in_g:
            in_g, g0 = True, y
        elif not b and in_g:
            in_g = False
            if y - g0 >= min_px:
                spans.append((g0, y))
    if in_g and len(is_blank) - g0 >= min_px:
        spans.append((g0, len(is_blank)))
    return spans


def detect_raw_strips(image: Image.Image,
                      white_thr: int = 240,
                      min_gap_px: int = 10,
                      min_block_px: int = 35) -> list:
    gray    = np.array(image.convert("L"))
    density = moving_avg(row_dark_density(gray, white_thr), w=9)
    blank   = density < 0.005
    strips, prev = [], 0
    for g0, g1 in gap_spans(blank, min_gap_px):
        if g0 - prev >= min_block_px:
            strips.append((prev, g0))
        prev = g1
    h = gray.shape[0]
    if h - prev >= min_block_px:
        strips.append((prev, h))
    return strips


def pair_notation_lyrics(strips: list, gray: np.ndarray,
                         white_thr: int = 240) -> list:
    if len(strips) < 2:
        return strips

    def dens(t, b):
        return (gray[t:b] < white_thr).mean() if b > t else 0.0

    def inter_gap(i):
        return strips[i + 1][0] - strips[i][1] if i + 1 < len(strips) else 99999

    out, i = [], 0
    while i < len(strips):
        t0, b0 = strips[i]
        if i + 1 < len(strips):
            t1, b1 = strips[i + 1]
            g_in   = inter_gap(i)
            g_out  = inter_gap(i + 1)
            d0, d1 = dens(t0, b0), dens(t1, b1)
            tight             = g_in < g_out * 0.60 and g_in < 110
            stave_above_lyric = g_in < 90 and d0 > d1 * 1.15
            if tight or stave_above_lyric:
                out.append((t0, b1))
                i += 2
                continue
        out.append(strips[i])
        i += 1
    return out


def group_singable(lines: list, n: int) -> list:
    groups = []
    for start in range(0, len(lines), n):
        chunk = lines[start: start + n]
        groups.append((chunk[0][0], chunk[-1][1]))
    return groups


def make_slide_crop(page_rgb, top, bottom, pad=22):
    w, h  = page_rgb.size
    t     = max(0, top - pad)
    b     = min(h, bottom + pad)
    crop  = page_rgb.crop((0, t, w, b))
    crop  = ImageOps.expand(crop, border=28, fill=(255, 255, 255))
    return crop


def add_image_slide(prs, img_path):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    sw, sh = prs.slide_width, prs.slide_height
    with Image.open(img_path) as im:
        iw, ih = im.size
    aspect = iw / ih
    max_w  = sw * 0.96
    max_h  = sh * 0.90
    fit_w  = max_w
    fit_h  = fit_w / aspect
    if fit_h > max_h:
        fit_h = max_h
        fit_w = fit_h * aspect
    left = (sw - fit_w) / 2
    top  = (sh - fit_h) / 2
    slide.shapes.add_picture(img_path, int(left), int(top), int(fit_w), int(fit_h))


# ══════════════════════════════════════════════════════════════════
#  Bible-slide helpers
# ══════════════════════════════════════════════════════════════════

def chunk_verses(verses: list, per_slide: int) -> list:
    """Split a list of verse strings into groups of per_slide."""
    return [verses[i: i + per_slide] for i in range(0, len(verses), per_slide)]


def build_bible_pptx(verses: list,
                     per_slide: int,
                     bg_hex: str,
                     text_hex: str,
                     ref_hex: str,
                     font_name: str,
                     font_size: int,
                     ref_size: int,
                     bold: bool,
                     align: str,
                     reference: str,
                     show_ref_each: bool,
                     line_spacing: float = 1.2) -> str:

    from pptx.util import Pt, Inches
    from pptx.oxml.ns import qn
    from lxml import etree

    prs = new_prs()
    sw  = prs.slide_width
    sh  = prs.slide_height
    pad = Inches(0.6)

    align_map = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}
    pp_align  = align_map.get(align, PP_ALIGN.CENTER)
    tr, tg, tb = hex_to_rgb(text_hex)
    rr, rg, rb = hex_to_rgb(ref_hex)

    groups = chunk_verses(verses, per_slide)
    total  = len(groups)

    def _add_text(slide, lines, left, top, width, height,
                  fname, fsize, fbold, fcolor_rgb, falign, spacing=1.2):
        txBox = slide.shapes.add_textbox(int(left), int(top), int(width), int(height))
        tf    = txBox.text_frame
        tf.word_wrap = True
        for i, line in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = falign
            pPr   = p._p.get_or_add_pPr()
            lSpc  = etree.SubElement(pPr, qn('a:lSpc'))
            sPct  = etree.SubElement(lSpc, qn('a:spcPct'))
            sPct.set('val', str(int(spacing * 100000)))
            run = p.add_run()
            run.text = line
            run.font.name  = fname
            run.font.size  = Pt(fsize)
            run.font.bold  = fbold
            run.font.color.rgb = RGBColor(*fcolor_rgb)

    for idx, group in enumerate(groups):
        slide = blank_slide(prs, bg_hex)

        _add_text(slide, group,
                  pad, pad,
                  sw - pad * 2, sh - pad * 2 - Inches(0.7),
                  font_name, font_size, bold, (tr, tg, tb), pp_align,
                  spacing=line_spacing)

        ref_label = reference if (show_ref_each or idx == total - 1) else ""
        if ref_label:
            _add_text(slide, [ref_label],
                      pad, sh - Inches(0.75),
                      sw - pad * 2, Inches(0.6),
                      font_name, ref_size, False, (rr, rg, rb),
                      PP_ALIGN.RIGHT)

    return save_prs(prs, f"bible_{uuid.uuid4()}.pptx")


# ══════════════════════════════════════════════════════════════════
#  Lyric-extract helpers
# ══════════════════════════════════════════════════════════════════

import re as _re

# Chord patterns: lines that are ONLY chord symbols like A, Bm, G/B, F#m7, Am/C, D7, Cadd9 …
# A lyric line may start with a capital but will have Khmer or lowercase letters after it.
_CHORD_LINE_RE = _re.compile(
    r'^[\s]*'                         # optional leading space
    r'(?:'
      r'[A-G]'                        # root note
      r'(?:b|#)?'                     # flat / sharp
      r'(?:m|maj|min|aug|dim|sus|add)?'   # quality
      r'(?:\d+)?'                     # extension (7, 9, 11 …)
      r'(?:/[A-G](?:b|#)?)?'         # slash chord  e.g. G/B
    r'[\s\-]*'                        # spacing between chords
    r')+'
    r'[\s]*$',                        # nothing else on the line
    _re.IGNORECASE
)

# Lines that look purely like music/notation markings (e.g.  "| | |", "||", bar numbers)
_NOTATION_LINE_RE = _re.compile(r'^[\s\|\-\.\d:]+$')

# Detect whether a line contains at least one Khmer Unicode character
_HAS_KHMER_RE = _re.compile(r'[\u1780-\u17FF\u19E0-\u19FF]')

# Detect whether a line contains at least some real alphabetic word (for English lyrics)
_HAS_WORD_RE = _re.compile(r'[a-zA-Z]{2,}')


def _is_chord_line(line: str) -> bool:
    """Return True if the entire line is just chord symbols — no real lyric content."""
    stripped = line.strip()
    if not stripped:
        return False
    # If the line contains Khmer, it is definitely a lyric line
    if _HAS_KHMER_RE.search(stripped):
        return False
    # Match pure chord line
    if _CHORD_LINE_RE.match(stripped):
        return True
    # Also catch short ALL-CAPS tokens that are likely chords (e.g. "Am  G  F  C")
    tokens = stripped.split()
    if all(_re.fullmatch(r'[A-G][#b]?(?:m|maj|min|aug|dim|sus|add)?(?:\d+)?(?:/[A-G][#b]?)?', t, _re.IGNORECASE)
           for t in tokens if t):
        return True
    return False


def _is_notation_line(line: str) -> bool:
    """Return True if the line is pure music notation marks (bars, dashes, numbers)."""
    return bool(_NOTATION_LINE_RE.match(line.strip())) and not _HAS_KHMER_RE.search(line)


def filter_lyrics(raw_text: str) -> str:
    """
    Remove chord lines, pure notation lines, and blank lines from OCR output.
    Keeps Khmer lines and English lyric lines (lines with real words).
    """
    kept = []
    for line in raw_text.splitlines():
        stripped = line.strip()
        if not stripped:
            continue
        if _is_chord_line(stripped):
            continue
        if _is_notation_line(stripped):
            continue
        # Keep if it has Khmer or has actual English words (>=2 alpha chars)
        if _HAS_KHMER_RE.search(stripped) or _HAS_WORD_RE.search(stripped):
            kept.append(stripped)
    return "\n".join(kept)


def _extract_lyric_strips_from_page(page_rgb: Image.Image) -> list:
    """
    Given a full page image that contains music notation rows alternating with
    Khmer lyric rows, return a list of PIL sub-images that contain ONLY the
    lyric text rows — NOT the staff/note-head rows.

    Strategy:
      1. Detect all horizontal content strips using a NARROW smoothing window
         (w=3) so dense stave rows don't bleed into adjacent lyric rows.
      2. Measure the ink density of each strip individually (no smoothing).
      3. Strips with density > STAVE_THR are notation rows → skip.
         Strips with density in (LYRIC_MIN, STAVE_THR] are lyric rows → keep.
      4. Very short strips (< 12 px) are slur/tie connectors → skip.

    Real-world typical values from scanned Khmer song sheets:
      • Stave rows:    ~20–50 % dark pixels  (staff lines + note heads)
      • Lyric rows:    ~2–12 % dark pixels   (Khmer glyphs, spaced out)
    """
    gray    = np.array(page_rgb.convert("L"))
    # Use narrow smoothing so stave/lyric boundaries stay sharp
    density = moving_avg(row_dark_density(gray, 240), w=3)
    blank   = density < 0.004   # < 0.4 % ink → truly blank row

    # Collect content strips
    strips, prev = [], 0
    for g0, g1 in gap_spans(blank, min_px=5):
        if g0 - prev >= 10:
            strips.append((prev, g0))
        prev = g1
    h = gray.shape[0]
    if h - prev >= 10:
        strips.append((prev, h))

    if not strips:
        return []

    # Measure each strip's actual ink density (no smoothing — raw band)
    def strip_density(t, b):
        return float((gray[t:b] < 240).mean()) if b > t else 0.0

    densities = [strip_density(t, b) for t, b in strips]

    # Adaptive threshold: median of all strip densities is our split point.
    # Notation rows sit above the median; lyric rows sit below.
    # We also enforce an absolute minimum (0.5 %) so pure-blank bands are excluded.
    if not densities:
        return []

    median_d = float(np.median(densities))
    # Use the midpoint between 0 and the median as the lyric upper bound
    # so that even if everything on the page is moderate density we still
    # pick the relatively sparser strips.
    STAVE_THR  = max(median_d * 0.75, 0.04)   # strips above this → notation
    LYRIC_MIN  = 0.004                          # strips below this → blank/skip
    MIN_HEIGHT = 12                             # px

    w_page = page_rgb.size[0]
    lyric_images = []
    for (top, bot), d in zip(strips, densities):
        if (bot - top) < MIN_HEIGHT:
            continue
        if d < LYRIC_MIN:
            continue
        if d >= STAVE_THR:
            continue   # notation / stave row

        # This is a lyric row
        t = max(0, top - 4)
        b = min(h, bot + 6)
        lyric_images.append(page_rgb.crop((0, t, w_page, b)))

    return lyric_images


def extract_text_from_image(img: Image.Image) -> str:
    try:
        import pytesseract
        # Explicitly point to Homebrew tesseract so it works in venv
        pytesseract.pytesseract.tesseract_cmd = "/opt/homebrew/bin/tesseract"
        # Try Khmer + English; fall back gracefully
        try:
            text = pytesseract.image_to_string(img, lang="khm+eng")
        except Exception:
            try:
                text = pytesseract.image_to_string(img, lang="eng")
            except Exception as e2:
                return f"[OCR error: {e2}]"
        return text.strip()
    except ImportError:
        return "[pytesseract not installed]"


def extract_lyrics_from_file(filepath: str, ext: str) -> str:
    """
    Extract ONLY lyric lines from a song PDF/image.
    For song sheets that have notation rows interleaved with lyric rows,
    we first visually isolate each lyric row before running OCR on it,
    so chord symbols and staff lines never reach the OCR engine.
    """
    import pytesseract
    pytesseract.pytesseract.tesseract_cmd = "/opt/homebrew/bin/tesseract"

    all_lines = []

    def ocr_strip(strip_img: Image.Image) -> str:
        """Run OCR on one lyric strip, return stripped text."""
        try:
            t = pytesseract.image_to_string(strip_img, lang="khm+eng")
        except Exception:
            try:
                t = pytesseract.image_to_string(strip_img, lang="eng")
            except Exception:
                return ""
        return t.strip()

    def process_page(page_rgb: Image.Image):
        # Try smart lyric-strip isolation first
        strips = _extract_lyric_strips_from_page(page_rgb)
        if strips:
            for strip in strips:
                t = ocr_strip(strip)
                if t:
                    # Each strip may produce multiple lines — add all non-empty ones
                    for ln in t.splitlines():
                        ln = ln.strip()
                        if ln:
                            all_lines.append(ln)
        else:
            # Fallback: OCR the whole page and rely on text filtering
            t = extract_text_from_image(page_rgb)
            if t:
                all_lines.append(t)

    if ext == ".pdf":
        try:
            pages = convert_from_path(filepath, dpi=250)
        except Exception as e:
            return f"[PDF conversion error: {e}]"
        for page in pages:
            process_page(page.convert("RGB"))
    else:
        try:
            img = Image.open(filepath).convert("RGB")
            process_page(img)
        except Exception as e:
            return f"[Image error: {e}]"

    if not all_lines:
        return "(no text detected)"

    # Still run the text-level chord filter as a safety net
    raw = "\n".join(all_lines)
    filtered = filter_lyrics(raw)
    return filtered if filtered.strip() else raw


def build_lyric_pptx(lyrics_text: str,
                     lines_per_slide: int,
                     bg_hex: str,
                     text_hex: str,
                     font_name: str,
                     font_size: int,
                     bold: bool,
                     align: str,
                     line_spacing: float = 1.5) -> str:
    """
    Build a lyric PPTX.
    line_spacing: multiple of line height (1.0 = single, 1.5 = 1.5x, 2.0 = double).
    Each line in a group gets its own paragraph so spacing is applied correctly.
    Font name is set on every run to ensure it survives in PowerPoint.
    """
    from pptx.util import Pt, Inches, Emu
    from pptx.oxml.ns import qn
    from lxml import etree

    prs = new_prs()
    sw  = prs.slide_width
    sh  = prs.slide_height
    pad = Inches(0.55)

    align_map = {
        "left":   PP_ALIGN.LEFT,
        "center": PP_ALIGN.CENTER,
        "right":  PP_ALIGN.RIGHT,
    }
    pp_align = align_map.get(align, PP_ALIGN.CENTER)
    r, g, b  = hex_to_rgb(text_hex)

    all_lines = [l for l in lyrics_text.splitlines() if l.strip()]
    if not all_lines:
        all_lines = ["(no lyrics found)"]

    groups = [all_lines[i: i + lines_per_slide]
              for i in range(0, len(all_lines), lines_per_slide)]

    for group in groups:
        slide  = blank_slide(prs, bg_hex)
        txBox  = slide.shapes.add_textbox(int(pad), int(pad),
                                          int(sw - pad * 2),
                                          int(sh - pad * 2))
        tf     = txBox.text_frame
        tf.word_wrap = True

        for line_idx, line_text in enumerate(group):
            # First line reuses the auto-created paragraph; others are new
            p = tf.paragraphs[0] if line_idx == 0 else tf.add_paragraph()
            p.alignment = pp_align

            # Line spacing via XML (pPr/lSpc)
            pPr = p._pPr if p._pPr is not None else p._p.get_or_add_pPr()
            lSpc = etree.SubElement(pPr, qn('a:lSpc'))
            spcPct = etree.SubElement(lSpc, qn('a:spcPct'))
            spcPct.set('val', str(int(line_spacing * 100000)))  # in 1/1000 %

            run = p.add_run()
            run.text = line_text
            run.font.name  = font_name
            run.font.size  = Pt(font_size)
            run.font.bold  = bold
            run.font.color.rgb = RGBColor(r, g, b)

    return save_prs(prs, f"lyrics_{uuid.uuid4()}.pptx")


# ══════════════════════════════════════════════════════════════════
#  Flask routes
# ══════════════════════════════════════════════════════════════════

@app.route("/")
def index():
    return render_template("dashboard.html")


# ── Notation cut ─────────────────────────────────────────────────
@app.route("/api/notation", methods=["POST"])
def api_notation():
    if "pdf" not in request.files:
        return jsonify({"error": "No file uploaded"}), 400

    file            = request.files["pdf"]
    lines_per_slide = int(request.form.get("lines_per_slide", 1))
    dpi             = int(request.form.get("dpi", 200))

    filename = str(uuid.uuid4()) + ".pdf"
    filepath = os.path.join(UPLOAD_FOLDER, filename)
    file.save(filepath)

    for f in os.listdir(OUTPUT_FOLDER):
        if f.endswith(".png"):
            os.remove(os.path.join(OUTPUT_FOLDER, f))

    try:
        pages = convert_from_path(filepath, dpi=dpi)
    except Exception as e:
        os.remove(filepath)
        return jsonify({"error": f"PDF conversion failed: {e}"}), 500

    prs = new_prs()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    for page in pages:
        page_rgb = page.convert("RGB")
        w, h     = page_rgb.size
        gray     = np.array(page_rgb.convert("L"))

        strips   = detect_raw_strips(page_rgb, white_thr=240,
                                     min_gap_px=10, min_block_px=35)
        if not strips:
            third  = h // 3
            strips = [(i * third, (i + 1) * third) for i in range(3)]

        singable = pair_notation_lyrics(strips, gray)
        groups   = group_singable(singable, lines_per_slide)

        for grp_top, grp_bot in groups:
            crop = make_slide_crop(page_rgb, grp_top, grp_bot, pad=22)
            if crop.size[1] < 40:
                continue
            img_path = os.path.join(OUTPUT_FOLDER, f"block_{uuid.uuid4()}.png")
            crop.save(img_path, "PNG")
            add_image_slide(prs, img_path)

    out_path = save_prs(prs, "notation_output.pptx")
    os.remove(filepath)

    return send_file(out_path, as_attachment=True,
                     download_name="notation_slides.pptx",
                     mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation")


# ── Bible slides ─────────────────────────────────────────────────
@app.route("/api/bible", methods=["POST"])
def api_bible():
    data          = request.get_json(force=True)
    raw_text      = data.get("text", "").strip()
    per_slide     = max(1, int(data.get("per_slide", 1)))
    bg_hex        = data.get("bg_color", "000000").lstrip("#")
    text_hex      = data.get("text_color", "FFFFFF").lstrip("#")
    ref_hex       = data.get("ref_color", "AAAAAA").lstrip("#")
    font_name     = data.get("font", "Arial")
    font_size     = max(10, int(data.get("font_size", 36)))
    ref_size      = max(10, int(data.get("ref_size", 20)))
    bold          = bool(data.get("bold", False))
    align         = data.get("align", "center")
    reference     = data.get("reference", "").strip()
    show_ref_each = bool(data.get("show_ref_each", False))
    line_spacing  = float(data.get("line_spacing", 1.2))

    if not raw_text:
        return jsonify({"error": "No verse text provided"}), 400

    # Split into verses — each non-empty line is one verse
    verses = [l.strip() for l in raw_text.splitlines() if l.strip()]

    out_path = build_bible_pptx(verses, per_slide, bg_hex, text_hex,
                                 ref_hex, font_name, font_size, ref_size,
                                 bold, align, reference, show_ref_each,
                                 line_spacing)

    return send_file(out_path, as_attachment=True,
                     download_name="bible_slides.pptx",
                     mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation")


# ── Extract lyrics ────────────────────────────────────────────────
@app.route("/api/extract-lyrics", methods=["POST"])
def api_extract_lyrics():
    """OCR the uploaded file and return extracted text as JSON."""
    if "file" not in request.files:
        return jsonify({"error": "No file uploaded"}), 400

    f    = request.files["file"]
    ext  = os.path.splitext(f.filename)[1].lower()
    if ext not in (".pdf", ".png", ".jpg", ".jpeg"):
        return jsonify({"error": "Only PDF, PNG, JPG files are supported"}), 400

    filename = str(uuid.uuid4()) + ext
    filepath = os.path.join(UPLOAD_FOLDER, filename)
    f.save(filepath)

    text = extract_lyrics_from_file(filepath, ext)
    os.remove(filepath)
    return jsonify({"text": text})


@app.route("/api/lyric-slides", methods=["POST"])
def api_lyric_slides():
    """Build a PPTX from provided lyric text + style settings."""
    data            = request.get_json(force=True)
    lyrics_text     = data.get("text", "").strip()
    lines_per_slide = max(1, int(data.get("lines_per_slide", 2)))
    bg_hex          = data.get("bg_color", "000000").lstrip("#")
    text_hex        = data.get("text_color", "FFFFFF").lstrip("#")
    font_name       = data.get("font", "Arial")
    font_size       = max(10, int(data.get("font_size", 40)))
    bold            = bool(data.get("bold", False))
    align           = data.get("align", "center")
    line_spacing    = float(data.get("line_spacing", 1.5))

    if not lyrics_text:
        return jsonify({"error": "No lyrics text provided"}), 400

    out_path = build_lyric_pptx(lyrics_text, lines_per_slide, bg_hex,
                                 text_hex, font_name, font_size, bold, align,
                                 line_spacing=line_spacing)

    return send_file(out_path, as_attachment=True,
                     download_name="lyric_slides.pptx",
                     mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation")


# ── Preview slides (returns JSON list of slide data for JS carousel) ──
@app.route("/api/preview-slides", methods=["POST"])
def api_preview_slides():
    """
    Returns a JSON array of slide objects for the live preview carousel.
    Each slide has: { bg, text, ref, font, fontSize, bold, align, refColor, refSize }
    """
    data      = request.get_json(force=True)
    kind      = data.get("kind", "bible")   # "bible" | "lyric"
    raw_text  = data.get("text", "").strip()
    per_slide = max(1, int(data.get("per_slide", 1)))

    if not raw_text:
        return jsonify([])

    lines  = [l for l in raw_text.splitlines() if l.strip()]
    groups = [lines[i: i + per_slide] for i in range(0, len(lines), per_slide)]

    slides = []
    ref      = data.get("reference", "")
    ref_each = bool(data.get("show_ref_each", False))
    total    = len(groups)

    for idx, group in enumerate(groups):
        show_ref = ref and (ref_each or idx == total - 1)
        slides.append({
            "bg":          data.get("bg_color", "#000000"),
            "text":        "\n".join(group),
            "ref":         ref if show_ref else "",
            "font":        data.get("font", "Arial"),
            "fontSize":    int(data.get("font_size", 36)),
            "bold":        bool(data.get("bold", False)),
            "align":       data.get("align", "center"),
            "color":       data.get("text_color", "#ffffff"),
            "refColor":    data.get("ref_color", "#aaaaaa"),
            "refSize":     int(data.get("ref_size", 20)),
            "lineSpacing": float(data.get("line_spacing", 1.5)),
        })
    return jsonify(slides)


# keep old /upload route alive for backward compat
@app.route("/upload", methods=["POST"])
def upload():
    return api_notation()


if __name__ == "__main__":
    app.run(debug=True, port=5000)


# ══════════════════════════════════════════════════════════════════
#  Image-analysis helpers
# ══════════════════════════════════════════════════════════════════

def row_dark_density(gray: np.ndarray, threshold: int = 240) -> np.ndarray:
    """Fraction of pixels darker than threshold in each row."""
    return (gray < threshold).mean(axis=1).astype(float)


def moving_avg(arr: np.ndarray, w: int = 9) -> np.ndarray:
    return np.convolve(arr, np.ones(w) / w, mode="same")


def gap_spans(is_blank: np.ndarray, min_px: int) -> list:
    """Contiguous True-runs that are at least min_px long."""
    spans, in_g, g0 = [], False, 0
    for y, b in enumerate(is_blank):
        if b and not in_g:
            in_g, g0 = True, y
        elif not b and in_g:
            in_g = False
            if y - g0 >= min_px:
                spans.append((g0, y))
    if in_g and len(is_blank) - g0 >= min_px:
        spans.append((g0, len(is_blank)))
    return spans


def detect_raw_strips(image: Image.Image,
                      white_thr: int = 240,
                      min_gap_px: int = 10,
                      min_block_px: int = 35) -> list:
    """
    Cut the page on every horizontal whitespace gap.
    Returns [(top, bottom), …] for every content strip,
    always using FULL page width.
    """
    gray    = np.array(image.convert("L"))
    density = moving_avg(row_dark_density(gray, white_thr), w=9)
    blank   = density < 0.005          # < 0.5 % ink → blank row

    strips, prev = [], 0
    for g0, g1 in gap_spans(blank, min_gap_px):
        if g0 - prev >= min_block_px:
            strips.append((prev, g0))
        prev = g1
    h = gray.shape[0]
    if h - prev >= min_block_px:
        strips.append((prev, h))
    return strips


def pair_notation_lyrics(strips: list, gray: np.ndarray,
                         white_thr: int = 240) -> list:
    """
    Each music line = [notation strip] immediately followed by [lyric strip].
    Merge them into one "singable line".

    Decision rule: two consecutive strips belong together when
      • the gap between them  < 90 px  AND
      • the first is noticeably denser than the second
        (staff lines / note-heads  vs  sparse Khmer glyphs), OR
      • the gap between them is clearly smaller than the gap after them
        (tight inner-pair, large separator between lines).
    """
    if len(strips) < 2:
        return strips

    def dens(t, b):
        return (gray[t:b] < white_thr).mean() if b > t else 0.0

    def inter_gap(i):
        return strips[i + 1][0] - strips[i][1] if i + 1 < len(strips) else 99999

    out, i = [], 0
    while i < len(strips):
        t0, b0 = strips[i]
        if i + 1 < len(strips):
            t1, b1  = strips[i + 1]
            g_in    = inter_gap(i)
            g_out   = inter_gap(i + 1)
            d0, d1  = dens(t0, b0), dens(t1, b1)

            tight   = g_in < g_out * 0.60 and g_in < 110
            stave_above_lyric = g_in < 90 and d0 > d1 * 1.15

            if tight or stave_above_lyric:
                out.append((t0, b1))
                i += 2
                continue
        out.append(strips[i])
        i += 1
    return out


def group_singable(lines: list, n: int) -> list:
    """
    Combine n singable lines into one slide group.
    top  = first line's top
    bottom = last line's bottom
    (whitespace between lines is preserved naturally).
    """
    groups = []
    for start in range(0, len(lines), n):
        chunk = lines[start : start + n]
        groups.append((chunk[0][0], chunk[-1][1]))
    return groups


def make_slide_crop(page_rgb: Image.Image,
                    top: int, bottom: int,
                    pad: int = 22) -> Image.Image:
    """
    Crop a horizontal band from the page (full width), add white padding,
    and return a clean PIL image ready to be placed on a slide.
    """
    w, h   = page_rgb.size
    t      = max(0,  top    - pad)
    b      = min(h,  bottom + pad)
    crop   = page_rgb.crop((0, t, w, b))
    crop   = ImageOps.expand(crop, border=28, fill=(255, 255, 255))
    return crop


def add_image_slide(prs: Presentation, img_path: str) -> None:
    """
    Add a blank white slide and place the image centred,
    scaled as large as possible (up to 96 % width / 90 % height).
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])   # blank layout

    # Explicit white background
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    sw, sh = prs.slide_width, prs.slide_height

    with Image.open(img_path) as im:
        iw, ih = im.size
    aspect = iw / ih

    max_w = sw * 0.96
    max_h = sh * 0.90
    fit_w = max_w
    fit_h = fit_w / aspect
    if fit_h > max_h:
        fit_h = max_h
        fit_w = fit_h * aspect

    left = (sw - fit_w) / 2
    top  = (sh - fit_h) / 2
    slide.shapes.add_picture(img_path, int(left), int(top),
                              int(fit_w), int(fit_h))



if __name__ == "__main__":
    app.run(debug=True, port=5000)
